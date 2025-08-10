from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO

def add_content_slide(prs, heading, bullet_points, image_url):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Heading box (top-left)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = heading
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Bullet points box (below heading)
    bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(4))
    bullets_frame = bullets_box.text_frame
    bullets_frame.word_wrap = True

    for point in bullet_points:
        p = bullets_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(60, 60, 60)

    # Image box (right side)
    if image_url:
        img_data = requests.get(image_url).content
        slide.shapes.add_picture(BytesIO(img_data), Inches(6), Inches(1.5), width=Inches(3.5))

def add_intro_or_conclusion_slide(prs, heading, text):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title at top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = heading
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Body text below
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(8), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.text = text
    body_frame.paragraphs[0].font.size = Pt(20)
    body_frame.paragraphs[0].font.color.rgb = RGBColor(60, 60, 60)

def create_ppt(content, image_fetcher):
    prs = Presentation()

    # Intro Slide
    add_intro_or_conclusion_slide(prs, content["title"], content["intro"])

    # Content Slides
    for slide in content["slides"]:
        heading = slide["heading"]
        points = slide["text"].split(". ")
        points = [p.strip() for p in points if p.strip()][:4]
        img_url = image_fetcher(heading)
        add_content_slide(prs, heading, points, img_url)

    # Conclusion Slide
    add_intro_or_conclusion_slide(prs, "Conclusion", content["conclusion"])

    return prs
